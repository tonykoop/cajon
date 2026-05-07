"""
Build Cajon-Family-Capstone.pptx — capstone slide deck.
Run: python build_pptx.py
Requires: pip install python-pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "Cajon-Family-Capstone.pptx"

INK = RGBColor(0x2B, 0x24, 0x17)
WARM = RGBColor(0x6B, 0x4F, 0x2C)
ACCENT = RGBColor(0xB9, 0x5A, 0x3B)
PAPER = RGBColor(0xFC, 0xF9, 0xF1)


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PAPER

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.6), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Cambria"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = INK

    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(8.6), Inches(1.0))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = "Cambria"
    p2.font.size = Pt(20)
    p2.font.italic = True
    p2.font.color.rgb = WARM


def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Cambria"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + b
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Cajón Family",
        "A parametric three-member box-drum family — Heifer Zephyr Instruments"
    )

    add_content_slide(prs, "Project intent", [
        "Document the cajón as a portfolio-grade parametric instrument, not as one box",
        "Three coordinated sizes: Compact (CJ-C), Standard (CJ-S), Bass (CJ-B)",
        "Four manufacturing variants per size: V1 finger-joint, V2 dovetail, V3 rabbet, V4 miter+spline",
        "Optional snare lineage: A none (Peruvian) / B guitar-string (flamenco) / C snare-wires",
        "First prototype: Standard, V1, snare option B"
    ])

    add_content_slide(prs, "Governing physics", [
        "Two co-dominant resonators couple to make the bass voice",
        "Front-plate (1,1) clamped rectangular bending mode — the 3 mm tapa",
        "Helmholtz cavity with back-port — closed walnut box + 4.5″ hole",
        "f_H = (c/2π)·√(A/(V·L_eff)),  L_eff = t_back + 1.7·(d/2)",
        "Standard prediction: plate (1,1) ≈ 84 Hz · cavity ≈ 96 Hz · coupled mode is the bass thump",
        "Risk A2: original workbook formula was missing the 1.7·r end correction"
    ])

    add_content_slide(prs, "Family table", [
        "Compact: 15 × 11 × 11 in, 4.0″ hole, predicted f_H ≈ 107 Hz",
        "Standard ⭐: 18 × 12 × 12 in, 4.5″ hole, predicted f_H ≈ 96 Hz",
        "Bass: 22 × 14 × 14 in, 5.0″ hole, predicted f_H ≈ 77 Hz",
        "Compact runs tightly coupled (plate ≈ cavity); Bass runs cavity-dominant",
        "Standard sits in the middle — first-prototype target"
    ])

    add_content_slide(prs, "Manufacturing variants", [
        "V1 Finger joint — table-saw jig OR CNC parametric. First-prototype default",
        "V2 Half-blind dovetail — CNC. Showpiece path, harder fixturing",
        "V3 Rabbet — table saw. Quick commercial baseline",
        "V4 Miter + spline — table saw + spline jig. Visual ties to inlay",
        "Same shell geometry, four interchangeable joinery paths",
        "All four variants share BOM except the joint-specific bit and feature stock"
    ])

    add_content_slide(prs, "CNC inlay strategy (Broinwood-derived)", [
        "Tapered ball-nose 12.4° bit — single bit cuts pocket AND insert",
        "Pocket = negative in walnut body panel; insert = positive in maple",
        "Tapered shoulder geometry creates press-fit with no visible glue line",
        "VCarve allowance setting is the single tuning knob",
        "Route INLAY BEFORE assembly — no clamping access on a closed box",
        "Keep inlay within central 70% of panel; 1″ keep-out from joints"
    ])

    add_content_slide(prs, "Build packet", [
        "design.md · bom.csv · sourcing.csv · cut-list.csv · validation.csv",
        "assembly-manual.md · supplier-rfq.md · drawing-brief.md · visual-bom-brief.md",
        "risks.md (17 entries, 5 categories) · wolfram-starter.wl",
        "drawings/ — dimensioned SVGs · cad/ — OpenSCAD master · cnc/ — bit + feed strategies",
        "site/ — public build-log site · scripts/ — xlsx + pptx + pdf generators"
    ])

    add_content_slide(prs, "Validation plan", [
        "Predicted f_H, plate (1,1), and coupled bass per member in validation.csv",
        "First-prototype acceptance: f_H ±15 Hz of prediction; plate (1,1) ±15%",
        "Slap-vs-bass spectrum delta ≥ +5 dB at 1 kHz",
        "Sustain T60 between 0.30–0.55 s",
        "Tapa screw-torque uniformity drives plate boundary symmetry",
        "If cents error > 50, revisit the family scaling law (not the build)"
    ])

    add_content_slide(prs, "Risks (highest-leverage)", [
        "A2: Helmholtz formula missing end-correction (corrected in design.md)",
        "S1: Tapa screw torque non-uniformity → asymmetric plate boundary",
        "SU1: 3 mm ply thickness inconsistency → plate (1,1) varies as h³",
        "F1: Inlay press-fit too tight cracks insert (test on scrap first)",
        "S4: Tapa fatigue at slap zone — long-term wear item, replaceable"
    ])

    add_content_slide(prs, "Next actions", [
        "Build the Standard prototype",
        "Measure f_H, plate (1,1), and the slap-vs-centre spectral delta",
        "Update validation.csv with measured data + cents error",
        "If error > 50 cents, revise family scaling law in workbook",
        "Apply corrected Helmholtz formula to cajon-design-table.xlsx",
        "Generate remaining drawings (sheets 02, 04, 06, 07, 08)",
        "Capture build-log photography per photo-shotlist.md"
    ])

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
